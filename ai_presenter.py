import os
import json
import logging
import subprocess
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PDFProcessor:
    """Módulo 6: Extractor de texto e imágenes vectorizadas de documentos PDF"""
    def __init__(self, output_dir: str = "assets/pdf_images"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
    def process(self, pdf_path: str) -> Tuple[str, List[str]]:
        logging.info("Módulo 6: Extrayendo texto y gráficas métricas del PDF aportado...")
        try:
            import fitz
            from PIL import Image
            from io import BytesIO
            
            doc = fitz.open(pdf_path)
            text_content = ""
            image_paths = []
            
            for page_index in range(len(doc)):
                page = doc[page_index]
                text_content += page.get_text() + "\n"
                
                # Usamos el motor de Layout para evitar imágenes fragmentadas o en mosaico (tiling)
                blocks = page.get_text("dict")["blocks"]
                for b_index, b in enumerate(blocks):
                    if b["type"] == 1:  # Bloque de tipo Imagen
                        bbox = b["bbox"] # (x0, y0, x1, y1)
                        width = bbox[2] - bbox[0]
                        height = bbox[3] - bbox[1]
                        
                        if width >= 150 and height >= 150:
                            try:
                                # Capturamos exactamente la caja limítrofe compaginada por el autor
                                rect = fitz.Rect(bbox)
                                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=rect)
                                
                                mode = "RGBA" if pix.alpha else "RGB"
                                pil_img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                                
                                if mode == "RGBA":
                                    bg = Image.new("RGB", pil_img.size, (255, 255, 255))
                                    bg.paste(pil_img, mask=pil_img.split()[-1])
                                    pil_img = bg
                                elif pil_img.mode != "RGB":
                                    pil_img = pil_img.convert("RGB")
                                
                                img_path = os.path.abspath(os.path.join(self.output_dir, f"pdf_img_p{page_index}_{b_index}.jpg"))
                                pil_img.save(img_path, "JPEG", quality=95)
                                image_paths.append(img_path)
                            except Exception as e:
                                pass
            
            # Limitar a ~80000 caracteres para estabilizar los LLMs convencionales
            return text_content[:80000], image_paths
        except Exception as e:
            logging.error(f"Error procesando PDF: {e}")
            return "", []

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class DeepResearchEngine:
    def __init__(self):
        self.api_key = os.getenv("OPENROUTER_API_KEY")
        self.model_name = os.getenv("OR_MODEL_CHOICE", "openai/gpt-4o-mini")

    def research_topic(self, topic: str, language: str) -> str:
        logging.info(f"Módulo 1: Deep Research sobre '{topic}' en idioma '{language}'...")
        if not self.api_key:
            return f"Contexto simulado sobre: {topic} para pruebas offline debido a que no ingresaste tu OpenRouter API Key."
        
        prompt = (f"Actúa como un experto investigador e historiador. Crea material robusto y profundo para armar una MASTERCLASS sobre '{topic}'. "
                  f"Incluye contexto histórico, hitos documentados y conclusiones analíticas con datos cuantitativos. "
                  f"TODA TU RESPUESTA DEBE ESTAR ESCRITA ABSOLUTAMENTE EN IDIOMA: {language.upper()}. NO USES OTRO IDIOMA.")
        try:
            from openai import OpenAI
            client = OpenAI(base_url="https://openrouter.ai/api/v1", api_key=self.api_key)
            response = client.chat.completions.create(
                model=self.model_name,
                messages=[{"role": "user", "content": prompt}],
                timeout=90
            )
            return response.choices[0].message.content
        except Exception as e:
            logging.error(f"Error OpenRouter DeepResearch: {e}")
            return f"Contexto offline fallback por falla de API: {topic}"

class ScriptGenerator:
    def __init__(self):
        self.api_key = os.getenv("OPENROUTER_API_KEY")
        self.model_name = os.getenv("OR_MODEL_CHOICE", "openai/gpt-4o-mini")

    def generate_slides(self, topic: str, research_context: str, num_slides: int, language: str) -> List[Dict[str, Any]]:
        logging.info(f"Módulo 2: Generando guion para {num_slides} diapositivas en {language}...")
        if not self.api_key:
            return [
                {"slide_number": i+1, "title": f"Sumario de {topic}" if i == 0 else f"Desarrollo {i}", 
                 "bullets": ["Punto 1", "Punto 2"], "narration_text": "Texto demo a narrar de prueba.", 
                 "image_prompt": "Abstract map graphic"}
                 for i in range(num_slides)
            ]

        prompt = f"""
        Diseña el guion para una MASTERCLASS PROFESIONAL de EXACTAMENTE {num_slides} diapositivas independientes sobre el tema "{topic}".
        TODO EL TEXTO (title, bullets, narration_text) DEBE ESTAR ESCRITO EN {language.upper()}.
        La diapositiva 1 debe ser el "Sumario/Índice" compilando el temario de las siguientes diapositivas.
        Devuelve SOLO UN ARRAY DE OBJETOS JSON puro, sin marcadores markdown, con este esquema exacto:
        [{{ "slide_number": int, "title": str, "bullets": [str, str, str], "narration_text": str, "image_prompt": str }}]
        
        INSTRUCCIONES CRÍTICAS:
        1. DEBES CREAR {num_slides} DIAPOSITIVAS. No te detengas hasta completar la número {num_slides} bajo ningún concepto.
        2. 'bullets': 3 a 5 puntos cortos y de fuerte impacto visual.
        3. 'narration_text': Guion MUY extenso, explicativo y profesional (estas serán las Notas del Orador para que las lea el humano).
        4. 'image_prompt': UNA SOLA PALABRA CLAVE SUSTANTIVA EN INGLES PURA (ej: "Iran", "President", "Oil") para localizar fotos correctas en Wikipedia orientadas a ilustrar cada diapositiva en concreto.
        
        Contexto Base:\n{research_context[:8000]}
        """
        try:
            from openai import OpenAI
            client = OpenAI(base_url="https://openrouter.ai/api/v1", api_key=self.api_key)
            response = client.chat.completions.create(
                model=self.model_name,
                messages=[{"role": "user", "content": prompt}],
                timeout=300,
                max_tokens=15000 # Solución al corte abrupto a las ~14 diapositivas
            )
            text = response.choices[0].message.content
            start = text.find('[')
            end = text.rfind(']')
            if start != -1 and end != -1:
                json_part = text[start:end+1]
                try:
                    return json.loads(json_part)
                except Exception as parse_e:
                    raise Exception(f"El LLM devolvió un código JSON mal formado o se ahogó generando. Trama final generada: {text[-200:]} | Error: {str(parse_e)}")
            else:
                logging.error("No se detectó un array JSON en la respuesta. Falló el LLM.")
                raise Exception("El modelo LLM no generó corchetes de array de datos. Se desvió de las instrucciones completas.")
        except Exception as e:
            logging.error(f"Error OpenRouter Guionización: {e}")
            raise Exception(f"Falla de API/LLM: {str(e)}")

class ImageGenerator:
    """Módulo 3 Original: API DALL-E 3"""
    def __init__(self, output_dir: str = "assets/img"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        self.openai_key = os.getenv("OPENAI_API_KEY")
        
    def generate_image(self, prompt: str, filename: str) -> Tuple[str, str]:
        filepath = os.path.abspath(os.path.join(self.output_dir, filename + ".jpg"))
        credit = "IA Generativa (DALL-E 3)"
        if not self.openai_key:
            return "", credit
        try:
            from openai import OpenAI
            from PIL import Image
            from io import BytesIO
            import requests

            client = OpenAI(api_key=self.openai_key)
            response = client.images.generate(model="dall-e-3", prompt=prompt, size="1024x1024", quality="standard", n=1)
            img_data = requests.get(response.data[0].url).content
            img = Image.open(BytesIO(img_data))
            if img.mode in ("RGBA", "P"): img = img.convert("RGB")
            img.save(filepath, "JPEG")
            return filepath, credit
        except Exception as e:
            return "", ""

class ImageSearcher:
    """Módulo 3 Alternativo: Buscador Wikipedia con Auto-Recorte y DuckDuckGo Fallback"""
    def __init__(self, output_dir: str = "assets/img"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
    def search_and_download(self, query: str, filename: str) -> Tuple[str, str]:
        filepath = os.path.abspath(os.path.join(self.output_dir, filename + ".jpg"))
        import requests
        from PIL import Image
        from io import BytesIO
        
        # User-Agent orgánico para engañar/evitar rate-limits severos de las APIs al servidor en Cloud
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36'}
        
        def save_img(content, credit_text):
            try:
                img = Image.open(BytesIO(content))
                if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info):
                    alpha = img.convert('RGBA').split()[-1]
                    bg = Image.new("RGB", img.size, (255, 255, 255))
                    bg.paste(img.convert('RGB'), mask=alpha)
                    img = bg
                elif img.mode != "RGB":
                    img = img.convert("RGB")
                img.save(filepath, "JPEG", quality=95)
                return filepath, credit_text
            except Exception as e:
                logging.error(f"Error en Cloud parseando JPG: {e}")
                return "", ""

        # 1er Vía: DuckDuckGo Images (Más tolerante a IPs de servidores AWS y Datacenters)
        try:
            from duckduckgo_search import DDGS
            results = DDGS().images(query, max_results=3)
            for res in results:
                url = res.get("image")
                if url:
                    try:
                        img_response = requests.get(url, headers=headers, timeout=5)
                        if img_response.status_code == 200:
                            path, credit_ret = save_img(img_response.content, f"Licencia Libre ({res.get('source', 'Web')})")
                            if path: return path, credit_ret
                    except:
                        continue
        except Exception as e:
            logging.warning(f"DDGS Fallback failed: {e}")
            
        # 2da Vía: Fallback Automático a la API de Wikimedia antigua 
        try:
            params = {
                'action': 'query', 'format': 'json', 'prop': 'pageimages',
                'generator': 'search', 'gsrsearch': query, 'gsrlimit': 3, 'pithumbsize': 1000
            }
            wiki_res = requests.get('https://en.wikipedia.org/w/api.php', params=params, headers=headers, timeout=10).json()
            if 'query' in wiki_res and 'pages' in wiki_res['query']:
                for page_id, page_data in wiki_res['query']['pages'].items():
                    if 'thumbnail' in page_data:
                        url = page_data['thumbnail']['source']
                        try:
                            img_response = requests.get(url, headers=headers, timeout=5)
                            if img_response.status_code == 200:
                                path, credit_ret = save_img(img_response.content, f"Wikimedia (Art.: {page_data.get('title', '')})")
                                if path: return path, credit_ret
                        except:
                            continue
        except Exception as e:
            logging.error(f"Error Wikipedia API: {e}")
            
        return "", ""

class TTSGenerator:
    """Módulo 5: Generador Neuronal de Voces Premium (Edge TTS) con control de velocidad (1.5x)"""
    def __init__(self, output_dir: str = "assets/audio"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
    def generate_audio(self, text: str, filename: str, voice: str = "es-ES-AlvaroNeural", speed: int = 0) -> str:
        filepath = os.path.abspath(os.path.join(self.output_dir, filename + ".mp3"))
        
        try:
            import asyncio
            import edge_tts
            
            async def run_tts():
                sign = "+" if speed >= 0 else ""
                rate_str = f"{sign}{speed}%"
                communicate = edge_tts.Communicate(text, voice, rate=rate_str)
                await communicate.save(filepath)
                
            asyncio.run(run_tts())
            logging.info(f"  [TTS] Autogenerado audio Neuronal ({voice}) al {speed}%: {filepath}")
            return filepath
        except Exception as e:
            logging.error(f"Error generando TTS Neuronal: {e}")
            return ""

class PresentationAssembler:
    """Módulo 4: Ensamblado con Diapositivas e Inyección de Notas del Orador"""
    def assemble(self, filename: str, slides_data: List[Dict[str, Any]], image_paths: Dict[int, Tuple[str, str]], audio_paths: Dict[int, str] = None, footer_text: str = "MBAI NATIVE"):
        if audio_paths is None: audio_paths = {}
        logging.info("Módulo 4: Ensamblando PPTX corporativo final...")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        COLOR_BG = RGBColor(15, 23, 42)
        COLOR_TITLE = RGBColor(248, 250, 252)
        COLOR_ACCENT = RGBColor(56, 189, 248)
        COLOR_TEXT = RGBColor(203, 213, 225)
        COLOR_MUTED = RGBColor(100, 116, 139)

        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = COLOR_BG
        
        main_topic = slides_data[0].get("title", "MBAI NATIVE") if slides_data else "MASTERCLASS"
        main_topic = main_topic.replace("Sumario de", "").replace("Sumario del", "").replace("Sumario", "MBAI Native - Origen").strip()
        
        tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = main_topic.upper()
        p.font.bold = True
        p.font.size = Pt(60)
        p.font.color.rgb = COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER
        
        sub = tf.add_paragraph()
        sub.text = "Una Masterclass generada por IA Autónoma"
        sub.font.size = Pt(28)
        sub.font.color.rgb = COLOR_ACCENT
        sub.alignment = PP_ALIGN.CENTER

        layouts = ["left_image", "right_image"]
        
        for index, slide_data in enumerate(slides_data):
            num = slide_data.get('slide_number', index+1)
            title = slide_data.get('title', f"Diapositiva {num}").upper()
            bullets = slide_data.get('bullets', [])
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG
            
            tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(1))
            tf_title = tb_title.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = title
            p_title.font.bold = True
            p_title.font.size = Pt(40)
            p_title.font.color.rgb = COLOR_TITLE
            
            line = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(3.0), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_ACCENT
            
            has_img = False
            if image_paths.get(num):
                img_path, credit = image_paths[num]
                if img_path and os.path.exists(img_path):
                    has_img = True
                    layout_mode = layouts[index % 2]
                    
                    from PIL import Image
                    try:
                        with Image.open(img_path) as px_img:
                            ratio = px_img.width / px_img.height
                    except:
                        ratio = 1.0
                        
                    is_wide = ratio > (5.0 / 4.5)
                    
                    if layout_mode == "right_image":
                        if is_wide:
                            pic = slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.8), width=Inches(5.0))
                        else:
                            pic = slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.8), height=Inches(4.5))
                        
                        tb_body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
                        cred_box = slide.shapes.add_textbox(Inches(7.5), Inches(6.4), Inches(5.0), Inches(0.5))
                    else:
                        if is_wide:
                            pic = slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), width=Inches(5.0))
                        else:
                            pic = slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), height=Inches(4.5))
                            
                        tb_body = slide.shapes.add_textbox(Inches(6.3), Inches(1.8), Inches(6.2), Inches(4.5))
                        cred_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(5.0), Inches(0.5))
                        
                    cred_p = cred_box.text_frame.paragraphs[0]
                    cred_p.text = f"Fuente: {credit}"
                    cred_p.font.size = Pt(11)
                    cred_p.font.italic = True
                    cred_p.font.color.rgb = COLOR_MUTED
            
            if not has_img:
                tb_body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.5))
            
            tf_body = tb_body.text_frame
            tf_body.word_wrap = True
            
            for i, bullet in enumerate(bullets):
                p = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
                p.text = f"• {bullet}"
                p.font.size = Pt(24) if has_img else Pt(30)
                p.font.color.rgb = COLOR_TEXT
                p.space_before = Pt(20)
                
            if audio_paths and audio_paths.get(num):
                audio_path = audio_paths.get(num)
                if os.path.exists(audio_path):
                    icon_path = os.path.abspath("assets/play_icon.jpg")
                    if not os.path.exists(icon_path):
                        if not os.path.exists("assets"): os.makedirs("assets")
                        try:
                            from PIL import Image, ImageDraw
                            img = Image.new('RGB', (100, 100), color=(56, 189, 248)) # Botón Celeste
                            d = ImageDraw.Draw(img)
                            d.polygon([(30, 20), (30, 80), (80, 50)], fill=(15, 23, 42)) # Triangulo Azul Oscuro
                            img.save(icon_path, quality=95)
                        except Exception:
                            pass
                    try:
                        # Ponemos el reproductor interactivo miniatura (50px) en la esquina superior derecha
                        slide.shapes.add_movie(audio_path, Inches(12.3), Inches(0.4), Inches(0.5), Inches(0.5), poster_frame_image=icon_path, mime_type='video/mp4')
                    except Exception as e:
                        logging.error(f"Error incrustando reproductor mp3 pptx: {e}")
                
            footer = slide.shapes.add_textbox(Inches(10.0), Inches(6.8), Inches(3.0), Inches(0.5))
            f_p = footer.text_frame.paragraphs[0]
            f_p.text = footer_text.upper()[:30] # Limite visual preventivo
            f_p.font.size = Pt(14)
            f_p.font.bold = True
            f_p.font.color.rgb = COLOR_ACCENT
            f_p.alignment = PP_ALIGN.RIGHT
            
            # INSERCIÓN CRÍTICA: Añadir todo el guion elaborado a las NOTAS DEL ORADOR de PowerPoint
            notas_del_orador = slide_data.get("narration_text", "")
            if notas_del_orador:
                # Automáticamente crea la vista de Notas para la Diapositiva e inyecta el texto
                slide.notes_slide.notes_text_frame.text = f"--- GUION PARA LA DIAPOSITIVA ---\n\n{notas_del_orador}"
                    
        filepath = os.path.abspath(filename)
        prs.save(filepath)
        logging.info(f"Presentación completada: {filepath}")
        return filepath

class AIPresenterPipeline:
    def __init__(self):
        self.researcher = DeepResearchEngine()
        self.script_gen = ScriptGenerator()
        self.img_gen = ImageGenerator()
        self.img_searcher = ImageSearcher()
        self.tts_gen = TTSGenerator()
        self.assembler = PresentationAssembler()
        self.pdf_processor = PDFProcessor()
        
    def run(self, topic: str, num_slides: int = 5, upload_gws: bool = True, image_source: str = "web", forced_name: str = None, language: str = "Español", generate_tts: bool = True, pdf_path: str = None, footer_text: str = "MBAI NATIVE", tts_speed: int = 0, tts_voice: str = "es-ES-AlvaroNeural"):
        logging.info(f"== PIPELINE TEMA: '{topic}' | SLIDES: {num_slides} | IDIOMA: {language} | IMGs: {image_source} | VOZ: {tts_voice} ==")
        
        pdf_image_paths = []
        if pdf_path and os.path.exists(pdf_path):
            logging.info("== MODO ESTRICTO: PDF DETECTADO ==")
            pdf_context, pdf_image_paths = self.pdf_processor.process(pdf_path)
            if pdf_context.strip():
                context = pdf_context
            else:
                context = self.researcher.research_topic(topic, language)
        else:
            context = self.researcher.research_topic(topic, language)
            
        slides = self.script_gen.generate_slides(topic, context, num_slides, language)
        if not slides:
            raise ValueError("Error crítico: el servidor backend no devolvió el array JSON final.")
            
        image_paths = {}
        audio_paths = {}
        
        pdf_img_index = 0
        for slide in slides:
            sn = slide.get('slide_number', 0)
            query = slide.get("image_prompt", "")
            if query:
                if image_source == "pdf":
                    if pdf_image_paths and pdf_img_index < len(pdf_image_paths):
                        image_paths[sn] = (pdf_image_paths[pdf_img_index], "Extracción Nativa del PDF Aportado")
                        pdf_img_index += 1
                elif image_source == "web":
                    path, credit = self.img_searcher.search_and_download(query, f"slide_{sn}")
                    if path: image_paths[sn] = (path, credit)
                else:
                    path, credit = self.img_gen.generate_image(query, f"slide_{sn}")
                    if path: image_paths[sn] = (path, credit)
            
            if generate_tts:
                narracion = slide.get("narration_text", "")
                if narracion:
                    a_path = self.tts_gen.generate_audio(narracion, f"narration_slide_{sn}", tts_voice, tts_speed)
                    if a_path: audio_paths[sn] = a_path
                    
        import re, time
        if forced_name:
            pptx_name = forced_name
        else:
            safe_topic = re.sub(r'[\\/*?:"<>|]', "", topic[:40])
            pptx_name = f"presentacion_{safe_topic.replace(' ', '_')}_{int(time.time())}.pptx"
            
        self.assembler.assemble(pptx_name, slides, image_paths, audio_paths, footer_text=footer_text)
        
        if upload_gws:
            logging.info("Subiendo presentación a Google Slides vía gws...")
            escaped_name = topic.replace('"', '\\"')
            cmd = f"gws drive files create --json '{{\"name\":\"{escaped_name}\",\"mimeType\":\"application/vnd.google-apps.presentation\"}}' --upload {pptx_name}"
            subprocess.run(cmd, shell=True, capture_output=True)

if __name__ == "__main__":
    pipeline = AIPresenterPipeline()
    pipeline.run("Test", num_slides=1, upload_gws=False, image_source="web", language="Español", generate_tts=False)
